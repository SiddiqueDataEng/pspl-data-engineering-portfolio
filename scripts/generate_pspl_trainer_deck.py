#!/usr/bin/env python3
"""
Generate a short trainer/overview PowerPoint for the PSPL lakehouse portfolio.

Requires: pip install python-pptx (see requirements.txt)

Usage (from repo root):
  python scripts/generate_pspl_trainer_deck.py

Output:
  docs/training/pspl_trainer_overview.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def _add_two_column_bullets(
    prs: Presentation,
    title: str,
    left_title: str,
    left: list[str],
    right_title: str,
    right: list[str],
) -> None:
    """Use blank layout and two text boxes."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True

    lx, ly, lw, lh = 0.5, 1.2, 4.5, 5.5
    rx, ry, rw, rh = 5.2, 1.2, 4.5, 5.5

    def fill_box(x, y, w, h, heading: str, lines: list[str]) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = heading
        p0.font.bold = True
        p0.font.size = Pt(22)
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)

    fill_box(lx, ly, lw, lh, left_title, left)
    fill_box(rx, ry, rw, rh, right_title, right)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "docs" / "training" / "pspl_trainer_overview.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "Pakistani social protection — data engineering portfolio",
        "Local lakehouse demo · Synthetic data · pspl.duckdb + Delta",
    )

    _add_bullets(
        prs,
        "Scope (in / out)",
        [
            "In: Bronze/Silver Delta, dbt marts in DuckDB, KPI SQL, Streamlit, tests, optional Airflow",
            "In: PSPL + refugee-style synthetic datasets for practice (not official stats)",
            "Out: Production PII, real programme numbers, streaming/CDC, cloud deploy of this exact layout",
            "Docs: docs/SCOPE_AND_CLOUD.md",
        ],
    )

    _add_bullets(
        prs,
        "Local pipeline order",
        [
            "0 Optional: datagenerator → data_large/",
            "1 ingest → delta_lake/bronze/",
            "2 Silver notebook → delta_lake/silver/",
            "3 dbt run / test → pspl.duckdb (repo root)",
            "4 sql-kpis, Streamlit dashboard; Airflow schedules the same steps in Docker",
        ],
    )

    _add_two_column_bullets(
        prs,
        "Cloud alternatives (selected)",
        "Local",
        [
            "PySpark + Delta on disk",
            "dbt-duckdb + pspl.duckdb",
            "Make / PowerShell scripts",
            "Streamlit on laptop",
        ],
        "Cloud-style",
        [
            "Databricks Spark + Delta on ADLS/S3",
            "dbt-databricks + Unity Catalog",
            "ADF / Workflows / managed Airflow",
            "Warehouse SQL + BI tools",
        ],
    )

    _add_bullets(
        prs,
        "Artifacts to verify",
        [
            "data_large/ — nine source files",
            "delta_lake/bronze and silver — Delta table folders",
            "pspl.duckdb — after dbt run (gitignored)",
            "Full command matrix: docs/RUN_EACH_COMPONENT.md",
        ],
    )

    _add_bullets(
        prs,
        "Trainer resources",
        [
            "docs/training/COMPLETE_TECHNICAL_TRAINER_GUIDE.md — labs and agendas",
            "docs/training/TRAINER_SLIDES_WITH_SPEAKER_NOTES.md — Marp-style manuscript",
            "Regenerate this deck: python scripts/generate_pspl_trainer_deck.py",
        ],
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
